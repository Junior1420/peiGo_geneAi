from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import numpy as np
import io
import datetime
import os

class PresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Configurar tema corporativo
        self.primary_color = RGBColor(0, 123, 255)  # Azul PeiGo
        self.secondary_color = RGBColor(40, 167, 69)  # Verde PeiGo
        self.accent_color = RGBColor(255, 193, 7)  # Amarillo PeiGo
        self.text_color = RGBColor(52, 58, 64)  # Gris oscuro

    def generate_presentation(self, user_state, conversation_history, output_path="PeiGo_Session_Report.pptx"):
        """Genera una presentación de PowerPoint basada en la sesión del usuario"""
        # Crear diapositiva de título
        self._create_title_slide("Informe de Sesión PeiGo", f"Generado: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}")

        # Crear diapositiva de resumen de interacción
        self._create_summary_slide(user_state, conversation_history)

        # Si hay simulaciones de crédito, crear diapositivas con detalles
        if user_state["credit_simulations"]:
            self._create_credit_simulation_slides(user_state["credit_simulations"])

        # Si se han discutido productos, crear diapositiva de productos
        if user_state["products_discussed"]:
            self._create_products_slide(user_state["products_discussed"])

        # Crear diapositiva de análisis de conversación
        self._create_conversation_analysis_slide(conversation_history)

        # Crear diapositiva final
        self._create_final_slide()

        # Guardar la presentación
        self.prs.save(output_path)
        return output_path

    def _create_title_slide(self, title, subtitle):
        """Crea la diapositiva de título"""
        slide_layout = self.prs.slide_layouts[0]  # Diseño de título
        slide = self.prs.slides.add_slide(slide_layout)

        # Configurar título
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = Pt(44)
        title_run.font.color.rgb = self.primary_color
        title_run.font.bold = True

        # Configurar subtítulo
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_para = subtitle_shape.text_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_para.runs[0]
        subtitle_run.font.size = Pt(24)
        subtitle_run.font.color.rgb = self.text_color

        # Agregar logo (en una implementación real se cargaría un archivo de imagen)
        # Para esta muestra, simplemente agregamos un texto como logo
        left = Inches(8.5)
        top = Inches(0.5)
        width = Inches(1)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "PeiGo"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.primary_color

    def _create_summary_slide(self, user_state, conversation_history):
        """Crea una diapositiva con el resumen de la interacción"""
        slide_layout = self.prs.slide_layouts[1]  # Diseño de título y contenido
        slide = self.prs.slides.add_slide(slide_layout)

        # Configurar título
        title_shape = slide.shapes.title
        title_shape.text = "Resumen de la Sesión"
        title_para = title_shape.text_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.primary_color

        # Configurar contenido
        content = slide.placeholders[1]
        tf = content.text_frame

        # Agregar información de la sesión
        p = tf.add_paragraph()
        p.text = "Datos de la interacción:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.secondary_color

        p = tf.add_paragraph()
        p.text = f"• Duración: {self._calculate_session_duration(conversation_history)} minutos"
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = f"• Mensajes intercambiados: {len(conversation_history)//2}"
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = f"• Productos consultados: {len(user_state['products_discussed'])}"
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = f"• Simulaciones de crédito: {len(user_state['credit_simulations'])}"
        p.font.size = Pt(14)

        # Agregar gráfico de interés (si hay simulaciones)
        if user_state["credit_simulations"]:
            sim = user_state["credit_simulations"][-1]  # Última simulación
            p = tf.add_paragraph()
            p.text = f"\nÚltima simulación: {sim['product']} por ${sim['amount']}"
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = self.secondary_color

            p = tf.add_paragraph()
            p.text = f"• Mensualidad: ${sim['monthly_payment']}"
            p.font.size = Pt(14)

            p = tf.add_paragraph()
            p.text = f"• Plazo: {sim['term']} meses"
            p.font.size = Pt(14)

            p = tf.add_paragraph()
            p.text = f"• Interés total: ${sim['total_interest']}"
            p.font.size = Pt(14)

    def _create_credit_simulation_slides(self, simulations):
        """Crea diapositivas con detalles de las simulaciones de crédito"""
        # Para cada simulación, crear una diapositiva
        for i, sim in enumerate(simulations):
            slide_layout = self.prs.slide_layouts[5]  # Diseño de título y contenido
            slide = self.prs.slides.add_slide(slide_layout)

            # Configurar título
            title_shape = slide.shapes.title
            title_shape.text = f"Simulación de {sim['product']}"
            title_para = title_shape.text_frame.paragraphs[0]
            title_run = title_para.runs[0]
            title_run.font.size = Pt(36)
            title_run.font.color.rgb = self.primary_color

            # Crear una tabla con los detalles de la simulación
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(1)
            rows, cols = 5, 2
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Establecer encabezados y datos
            table.cell(0, 0).text = "Parámetro"
            table.cell(0, 1).text = "Valor"
            table.cell(1, 0).text = "Monto del crédito"
            table.cell(1, 1).text = f"${sim['amount']}"
            table.cell(2, 0).text = "Plazo"
            table.cell(2, 1).text = f"{sim['term']} meses"
            table.cell(3, 0).text = "Pago mensual"
            table.cell(3, 1).text = f"${sim['monthly_payment']}"
            table.cell(4, 0).text = "Total a pagar"
            table.cell(4, 1).text = f"${sim['total_payment']}"

            # Dar formato a la tabla
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row, col)
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(14)
                    if row == 0:  # Encabezados
                        para.font.bold = True
                        para.font.color.rgb = self.primary_color

            # Agregar gráfico de distribución (capital vs interés)
            self._add_loan_distribution_chart(slide, sim)

    def _add_loan_distribution_chart(self, slide, simulation):
        """Agrega un gráfico de distribución de capital e interés"""
        chart_data = CategoryChartData()
        chart_data.categories = ['Capital', 'Interés']
        
        # Calcular capital e interés
        capital = simulation['amount']
        interest = simulation['total_interest']
        chart_data.add_series('Distribución', (capital, interest))

        # Crear gráfico
        x, y, cx, cy = Inches(1), Inches(4.5), Inches(8), Inches(3)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        
        # Personalizar gráfico
        chart.has_legend = True
        chart.legend.position = 2  # Colocar leyenda a la derecha
        chart.legend.include_in_layout = False
        
        # Agregar etiquetas
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '$#,##0'
        data_labels.font.size = Pt(12)
        data_labels.font.bold = True

    def _create_products_slide(self, products_discussed):
        """Crea una diapositiva con información sobre los productos discutidos"""
        slide_layout = self.prs.slide_layouts[1]  # Diseño de título y contenido
        slide = self.prs.slides.add_slide(slide_layout)

        # Configurar título
        title_shape = slide.shapes.title
        title_shape.text = "Productos de Interés"
        title_para = title_shape.text_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.primary_color

        # Configurar contenido
        content = slide.placeholders[1]
        tf = content.text_frame

        for product_name in products_discussed:
            p = tf.add_paragraph()
            p.text = product_name
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = self.secondary_color
            
            # En una implementación real, buscaríamos la información detallada
            # de cada producto en la base de conocimiento
            p = tf.add_paragraph()
            p.text = "• Descripción: Producto financiero de PeiGo diseñado para satisfacer necesidades específicas."
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = "• Beneficios: Tasas competitivas, aprobación rápida, atención personalizada."
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = "• Requisitos: Identificación oficial, comprobante de ingresos, buen historial crediticio."
            p.font.size = Pt(14)
            
            # Espacio entre productos
            p = tf.add_paragraph()
            p.text = ""

    def _create_conversation_analysis_slide(self, conversation_history):
        """Crea una diapositiva con análisis de la conversación"""
        slide_layout = self.prs.slide_layouts[1]  # Diseño de título y contenido
        slide = self.prs.slides.add_slide(slide_layout)

        # Configurar título
        title_shape = slide.shapes.title
        title_shape.text = "Análisis de la Conversación"
        title_para = title_shape.text_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.primary_color

        # Configurar contenido
        content = slide.placeholders[1]
        tf = content.text_frame

        # Análisis simple de la conversación
        # En una implementación real, se haría un análisis más sofisticado
        p = tf.add_paragraph()
        p.text = "Puntos clave de la conversación:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.secondary_color

        # Extraer temas principales (simplificado)
        topics = self._extract_conversation_topics(conversation_history)
        for topic, count in topics.items():
            p = tf.add_paragraph()
            p.text = f"• {topic.capitalize()}: mencionado {count} veces"
            p.font.size = Pt(14)

        # Agregar citas destacadas
        p = tf.add_paragraph()
        p.text = "\nPreguntas destacadas del cliente:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.secondary_color

        # Obtener algunas preguntas del cliente (simplificado)
        questions = self._extract_user_questions(conversation_history)
        for question in questions[:3]:  # Mostrar hasta 3 preguntas
            p = tf.add_paragraph()
            p.text = f'• "{question}"'
            p.font.size = Pt(14)
            p.font.italic = True

    def _create_final_slide(self):
        """Crea la diapositiva final con información de contacto"""
        slide_layout = self.prs.slide_layouts[6]  # Diseño en blanco
        slide = self.prs.slides.add_slide(slide_layout)

        # Agregar título
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = "¡Gracias por considerar a PeiGo!"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.primary_color

        # Agregar información de contacto
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)
        contact_box = slide.shapes.add_textbox(left, top, width, height)
        tf = contact_box.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = "Información de Contacto:"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.secondary_color

        p = tf.add_paragraph()
        p.text = "Teléfono: 01-800-PEIGO-01"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)

        p = tf.add_paragraph()
        p.text = "Email: contacto@peigo.com"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)

        p = tf.add_paragraph()
        p.text = "Sitio web: www.peigo.com"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)

    def _calculate_session_duration(self, conversation_history):
        """Calcula la duración aproximada de la sesión en minutos"""
        if len(conversation_history) < 2:
            return 0
            
        try:
            start_time = datetime.datetime.fromisoformat(conversation_history[0]["timestamp"])
            end_time = datetime.datetime.fromisoformat(conversation_history[-1]["timestamp"])
            duration = (end_time - start_time).total_seconds() / 60
            return round(duration, 1)
        except:
            return 5  # Valor predeterminado si hay error

    def _extract_conversation_topics(self, conversation_history):
        """Extrae temas principales de la conversación (simplificado)"""
        topics = {
            "crédito": 0,
            "productos": 0,
            "promociones": 0,
            "requisitos": 0,
            "pagos": 0
        }
        
        for message in conversation_history:
            content = message["content"].lower()
            for topic in topics:
                if topic in content:
                    topics[topic] += 1
                    
        # Filtrar temas no mencionados
        return {topic: count for topic, count in topics.items() if count > 0}

    def _extract_user_questions(self, conversation_history):
        """Extrae preguntas del usuario (simplificado)"""
        questions = []
        
        for message in conversation_history:
            if message["role"] == "user" and "?" in message["content"]:
                content = message["content"]
                # Extraer oraciones que terminen con interrogación
                for sentence in content.split("."):
                    if "?" in sentence:
                        question = sentence.strip()
                        if question and len(question) > 10:  # Filtrar preguntas demasiado cortas
                            questions.append(question)
                            
        return questions


# Ejemplo de uso
if __name__ == "__main__":
    # Datos de ejemplo
    user_state = {
        "name": "Juan Pérez",
        "identified": True,
        "interests": ["crédito personal", "bajo interés"],
        "credit_simulations": [
            {
                "timestamp": "2025-03-15T14:30:45",
                "product": "Crédito Personal",
                "amount": 25000,
                "term": 12,
                "monthly_payment": 2291.67,
                "total_payment": 27500.04,
                "total_interest": 2500.04
            }
        ],
        "products_discussed": ["PeiGo Flex", "PeiGo Start"]
    }
    
    conversation_history = [
        {"role": "user", "content": "Hola, ¿qué es PeiGo?", "timestamp": "2025-03-15T14:25:10"},
        {"role": "assistant", "content": "¡Hola! PeiGo es una plataforma financiera...", "timestamp": "2025-03-15T14:25:15"},
        {"role": "user", "content": "¿Qué productos ofrecen?", "timestamp": "2025-03-15T14:26:30"},
        {"role": "assistant", "content": "En PeiGo ofrecemos varios productos...", "timestamp": "2025-03-15T14:26:45"}
    ]
    
    # Generar presentación
    generator = PresentationGenerator()
    output_path = generator.generate_presentation(user_state, conversation_history)
    print(f"Presentación generada y guardada como: {output_path}")
